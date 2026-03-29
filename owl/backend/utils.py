import fitz # PyMuPDF
import logging
import os
from pptx import Presentation
from docx import Document

logger = logging.getLogger("owl-backend")

def extract_text_from_pdf(pdf_path: str) -> str:
    """Extracts text from a PDF file using PyMuPDF."""
    try:
        doc = fitz.open(pdf_path)
        text = ""
        for page in doc:
            text += page.get_text()
        doc.close()
        return text.strip()
    except Exception as e:
        logger.error(f"Failed to extract text from PDF: {str(e)}")
        raise e

def extract_text_from_pptx(pptx_path: str) -> str:
    """Extracts text from a PPTX file using python-pptx."""
    try:
        prs = Presentation(pptx_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text).strip()
    except Exception as e:
        logger.error(f"Failed to extract text from PPTX: {str(e)}")
        raise e

def extract_text_from_docx(docx_path: str) -> str:
    """Extracts text from a DOCX file using python-docx."""
    try:
        doc = Document(docx_path)
        return "\n".join([para.text for para in doc.paragraphs]).strip()
    except Exception as e:
        logger.error(f"Failed to extract text from DOCX: {str(e)}")
        raise e

def extract_text_any(file_path: str) -> str:
    """Detects file type based on extension and extracts text."""
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    elif ext in [".docx", ".doc"]:
        return extract_text_from_docx(file_path)
    elif ext == ".txt":
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read().strip()
    else:
        logger.error(f"Unsupported file format: {ext}")
        raise ValueError(f"Format non supporté: {ext}")
